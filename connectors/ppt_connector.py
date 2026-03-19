"""
PPT / PPTX Connector
--------------------
Extracts all text from every slide of a PowerPoint file using python-pptx.
Returns a list of LangChain Document objects, one per slide.
"""
from __future__ import annotations

import sys
from pathlib import Path
from typing import List

from langchain_core.documents import Document
from pptx import Presentation
from pptx.util import Pt

from logger.custom_logger import CustomLogger
from exception.custom_exception import DocumentPortalException

log = CustomLogger().get_logger(__name__)


def load_ppt(path: Path) -> List[Document]:
    """
    Load a .ppt or .pptx file and return one Document per slide.

    Args:
        path: Absolute path to the .pptx file.

    Returns:
        List of Document objects with slide text and metadata.
    """
    log.info("PPT connector: starting load", file=str(path))
    try:
        prs = Presentation(str(path))
        docs: List[Document] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text_parts: List[str] = []

            for shape in slide.shapes:
                # Extract text from text frames
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            slide_text_parts.append(line)

                # Extract text from tables inside slides
                if shape.has_table:
                    log.debug("PPT connector: found table in slide", slide=slide_num)
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip(" |"):
                            slide_text_parts.append(row_text)

            slide_text = "\n".join(slide_text_parts)

            if not slide_text.strip():
                log.warning(
                    "PPT connector: slide has no extractable text — skipping",
                    file=str(path),
                    slide=slide_num,
                )
                continue

            docs.append(
                Document(
                    page_content=slide_text,
                    metadata={
                        "source": str(path),
                        "file_type": path.suffix.lower(),
                        "slide_number": slide_num,
                    },
                )
            )

        log.info(
            "PPT connector: load complete",
            file=str(path),
            slides_loaded=len(docs),
            total_slides=len(prs.slides),
        )
        return docs

    except Exception as e:
        log.error("PPT connector: failed to load file", file=str(path), error=str(e))
        raise DocumentPortalException(f"Failed to load PPT file: {path.name}", e) from e
