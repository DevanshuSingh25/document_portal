from __future__ import annotations
import sqlite3
import pandas as pd
from pathlib import Path
from typing import Iterable, List

from pptx import Presentation
from langchain_core.documents import Document
from langchain_community.document_loaders import (
    PyPDFLoader, Docx2txtLoader, TextLoader,
    UnstructuredMarkdownLoader,
)

from logger.custom_logger import CustomLogger
from exception.custom_exception import DocumentPortalException

log = CustomLogger().get_logger(__name__)

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".txt", ".md",
    ".ppt", ".pptx",
    ".xlsx", ".csv",
    ".db", ".sqlite",
}


# ── Inline loaders ──────────────────────────────────────────────

def _load_pptx(path: Path) -> List[Document]:
    docs = []
    for i, slide in enumerate(Presentation(str(path)).slides):
        text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text")).strip()
        if text:
            docs.append(Document(page_content=text, metadata={"source": str(path), "slide": i + 1}))
    log.info("PPT loaded", file=path.name, slides=len(docs))
    return docs


def _load_xlsx(path: Path) -> List[Document]:
    docs = []
    for sheet, df in pd.read_excel(path, sheet_name=None).items():
        content = df.to_string(index=False)
        if content.strip():
            docs.append(Document(page_content=content, metadata={"source": str(path), "sheet": sheet}))
    log.info("XLSX loaded", file=path.name, sheets=len(docs))
    return docs


def _load_csv(path: Path) -> List[Document]:
    content = pd.read_csv(path).to_string(index=False)
    log.info("CSV loaded", file=path.name)
    return [Document(page_content=content, metadata={"source": str(path)})]


def _load_sqlite(path: Path) -> List[Document]:
    docs = []
    conn = sqlite3.connect(str(path))
    tables = conn.execute("SELECT name FROM sqlite_master WHERE type='table';").fetchall()
    for (table,) in tables:
        df = pd.read_sql_query(f'SELECT * FROM "{table}"', conn)
        if not df.empty:
            docs.append(Document(
                page_content=df.to_string(index=False),
                metadata={"source": str(path), "table": table},
            ))
    conn.close()
    log.info("SQLite loaded", file=path.name, tables=len(docs))
    return docs


# ── Main loader ─────────────────────────────────────────────────

def load_documents(paths: Iterable[Path]) -> List[Document]:
    docs: List[Document] = []
    try:
        for p in paths:
            ext = p.suffix.lower()
            if ext == ".pdf":
                docs.extend(PyPDFLoader(str(p)).load())
            elif ext == ".docx":
                docs.extend(Docx2txtLoader(str(p)).load())
            elif ext == ".txt":
                docs.extend(TextLoader(str(p), encoding="utf-8").load())
            elif ext == ".md":
                docs.extend(UnstructuredMarkdownLoader(str(p)).load())
            elif ext in {".ppt", ".pptx"}:
                docs.extend(_load_pptx(p))
            elif ext == ".xlsx":
                docs.extend(_load_xlsx(p))
            elif ext == ".csv":
                docs.extend(_load_csv(p))
            elif ext in {".db", ".sqlite"}:
                docs.extend(_load_sqlite(p))
            else:
                log.warning("Unsupported extension skipped", file=str(p))
                continue
            log.info("File loaded", file=p.name)
        log.info("Documents loaded", count=len(docs))
        return docs
    except DocumentPortalException:
        raise
    except Exception as e:
        log.error("Failed loading documents", error=str(e))
        raise DocumentPortalException("Error loading documents", e) from e


# ── Concatenation helpers ────────────────────────────────────────

def concat_for_analysis(docs: List[Document]) -> str:
    parts = []
    for d in docs:
        src = d.metadata.get("source") or d.metadata.get("file_path") or "unknown"
        parts.append(f"\n--- SOURCE: {src} ---\n{d.page_content}")
    return "\n".join(parts)


def concat_for_comparison(ref_docs: List[Document], act_docs: List[Document]) -> str:
    return f"<<REFERENCE_DOCUMENTS>>\n{concat_for_analysis(ref_docs)}\n\n<<ACTUAL_DOCUMENTS>>\n{concat_for_analysis(act_docs)}"